import tkinter as tk
from tkinter import messagebox
from tkinter import filedialog
from bs4 import BeautifulSoup
from ttkbootstrap import ttk, Style
import sys, os
import img2pdf
import re
import requests
from time import localtime, strftime
from os import listdir, walk
from os.path import isfile, join
import threading
import shutil
from pptx import Presentation
from pptx.util import Inches

CURRENT = os.path.dirname(__file__)

class SlideShareDownloader:
    def __init__(self, root):
        self.root = root
        self.root.title("SlideShare Downloader")
        self.root.geometry("500x400")

        self.download_folder = None
        self.output_file = None

        self.url_label = ttk.Label(self.root, text="Enter SlideShare URL:")
        self.url_label.pack(pady=10)

        self.url_entry = ttk.Entry(self.root, width=50)
        self.url_entry.pack(pady=5)

        self.format_label = ttk.Label(self.root, text="Select Output Format:")
        self.format_label.pack(pady=5)

        self.format_var = tk.StringVar(value="PDF")
        self.format_dropdown = ttk.Combobox(self.root, textvariable=self.format_var, values=["PDF", "PPT"])
        self.format_dropdown.pack(pady=5)

        self.download_button = ttk.Button(self.root, text="Download", command=self.start_download)
        self.download_button.pack(pady=20)

        # Circular Progress Bar (indeterminate mode)
        self.progress = ttk.Progressbar(self.root, orient="horizontal", length=300, mode="indeterminate")
        self.progress.pack(pady=10)

        self.open_folder_button = ttk.Button(self.root, text="Open Folder", state=tk.DISABLED, command=self.open_folder)
        self.open_folder_button.pack(pady=5)

        self.open_file_button = ttk.Button(self.root, text="Open Downloaded File", state=tk.DISABLED,
                                           command=self.open_downloaded_file)
        self.open_file_button.pack(pady=5)

        self.clear_button = ttk.Button(self.root, text="Clear", command=self.clear_input)
        self.clear_button.pack(pady=5)

        self.author_label = ttk.Label(self.root, text="Author - Kunal Saraswat | Project No. 5 | Version FINAL",
                                      font=('Helvetica', 10))
        self.author_label.pack(side=tk.BOTTOM, pady=10)

    def clear_input(self):
        self.url_entry.delete(0, tk.END)
        self.open_folder_button.config(state=tk.DISABLED)
        self.open_file_button.config(state=tk.DISABLED)
        self.progress.stop()

    def start_download(self):
        url = self.url_entry.get().strip()
        if not url:
            messagebox.showerror("Error", "Please enter a SlideShare URL.")
            return

        if not url.startswith(("http://", "https://")):
            url = "https://" + url

        output_format = self.format_var.get()
        threading.Thread(target=self.download_images, args=(url, output_format), daemon=True).start()
        self.progress.start()

    def download_images(self, url, output_format):
        try:
            html = requests.get(url).content
            soup = BeautifulSoup(html, "html.parser")

            title = "".join((CURRENT + "/pdf_images", strftime("/%Y%m%d_%H%M%S", localtime())))

            images = soup.find_all("img", {"data-testid": "vertical-slide-image"})
            image_url = ""

            for image in images:
                image_url = image.get("srcset").split("w, ")[-1].split(" ")[0]
                if image_url.endswith(".jpg"):
                    break

            i = 1
            image_url_prefix = image_url.rstrip("-2048.jpg")
            image_url_prefix = image_url_prefix.rstrip("0123456789")
            image_url_prefix = image_url_prefix.rstrip("-")
            output_file = re.sub("[^0-9a-zA-Z]+", "_", image_url_prefix.split("/")[-1])
            output_file += "." + output_format.lower()

            for image in images:
                image_url = image_url_prefix + "-" + str(i) + "-2048.jpg"
                r = requests.get(image_url)

                if not os.path.exists(title):
                    os.makedirs(title)

                filename = str(i) + ".jpg"
                i += 1

                with open(title + "/" + filename, "wb") as f:
                    f.write(r.content)

            if output_format == "PDF":
                self.convert_pdf(title, output_file)
            elif output_format == "PPT":
                self.convert_ppt(title, output_file)

        except Exception as e:
            messagebox.showerror("Error", f"An error occurred while downloading: {str(e)}")
        finally:
            self.progress.stop()

    def convert_pdf(self, img_dir_name, pdf_f):
        try:
            f = []
            for dirpath, dirnames, filenames in walk(join(CURRENT, img_dir_name)):
                f.extend(filenames)
                break
            f = ["%s/%s" % (img_dir_name, x) for x in f]

            def atoi(text):
                return int(text) if text.isdigit() else text

            def natural_keys(text):
                return [atoi(c) for c in re.split(r"(\d+)", text)]

            f.sort(key=natural_keys)

            pdf_bytes = img2pdf.convert(f, dpi=300, x=None, y=None)
            doc = open(pdf_f, "wb")
            doc.write(pdf_bytes)
            doc.close()

            if os.path.exists(img_dir_name):
                shutil.rmtree(img_dir_name)

            self.output_file = pdf_f
            self.download_folder = os.path.dirname(pdf_f)

            messagebox.showinfo("Success", f"Download and conversion complete!\nSaved as: {pdf_f}")
            self.open_folder_button.config(state=tk.NORMAL)
            self.open_file_button.config(state=tk.NORMAL)

        except Exception as e:
            messagebox.showerror("Error", f"An error occurred while converting to PDF: {str(e)}")

    def convert_ppt(self, img_dir_name, ppt_f):
        try:
            prs = Presentation()
            f = []
            for dirpath, dirnames, filenames in walk(join(CURRENT, img_dir_name)):
                f.extend(filenames)
                break
            f = ["%s/%s" % (img_dir_name, x) for x in f]

            def atoi(text):
                return int(text) if text.isdigit() else text

            def natural_keys(text):
                return [atoi(c) for c in re.split(r"(\d+)", text)]

            f.sort(key=natural_keys)

            for img_path in f:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(10))

            prs.save(ppt_f)

            if os.path.exists(img_dir_name):
                shutil.rmtree(img_dir_name)

            self.output_file = ppt_f
            self.download_folder = os.path.dirname(ppt_f)

            messagebox.showinfo("Success", f"Download and conversion complete!\nSaved as: {ppt_f}")
            self.open_folder_button.config(state=tk.NORMAL)
            self.open_file_button.config(state=tk.NORMAL)

        except Exception as e:
            messagebox.showerror("Error", f"An error occurred while converting to PPT: {str(e)}")

    def open_folder(self):
        if self.download_folder:
            os.startfile(self.download_folder)

    def open_downloaded_file(self):
        if self.output_file:
            os.startfile(self.output_file)

if __name__ == "__main__":
    root = tk.Tk()
    style = Style(theme="cyborg")
    app = SlideShareDownloader(root)
    root.mainloop()